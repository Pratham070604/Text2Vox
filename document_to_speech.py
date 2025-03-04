import os
import edge_tts
import PyPDF2
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
import mysql.connector

# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host="localhost",
        user="root",
        password="root",
        database="user_db"
    )

# Save document in database
def save_document(name, file_path):
    conn = get_db_connection()
    cursor = conn.cursor()

    query = "INSERT INTO documents (name, file_path) VALUES (%s, %s)"
    cursor.execute(query, (name, file_path))
    
    document_id = cursor.lastrowid  # Get the inserted document ID
    conn.commit()
    
    cursor.close()
    conn.close()
    return document_id  # Return document ID to link with audio


# Save generated audio in `document_audio` table
def save_document_audio(document_id, model, file_path):
    conn = get_db_connection()
    cursor = conn.cursor()

    query = "INSERT INTO document_audio (document_id, model, file_path) VALUES (%s, %s, %s)"
    cursor.execute(query, (document_id, model, file_path))
    
    conn.commit()
    cursor.close()
    conn.close()
    
# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text

# Function to extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from TXT
def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

# Function to extract text from HTML
def extract_text_from_html(html_path):
    with open(html_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html.parser')
        return soup.get_text()

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

# Function to convert text to speech using edge-tts
async def text_to_speech_edge_tts(text, output_file_path, voice):
    communicate = edge_tts.Communicate(text, voice=voice)  # Use the selected voice
    await communicate.save(output_file_path)
    print(f"Speech saved as {output_file_path}")

# Main function to process various documents and convert text to speech
# Process document and convert to speech
async def document_to_speech(doc_path, selected_voice):
    file_name = os.path.basename(doc_path)
    
    # ✅ Save document info to DB first
    document_id = save_document(file_name, doc_path)

    # Determine file type
    file_type = doc_path.rsplit('.', 1)[1].lower()

    # Extract text based on file type
    if file_type == 'pdf':
        text = extract_text_from_pdf(doc_path)
    elif file_type == 'docx':
        text = extract_text_from_docx(doc_path)
    elif file_type == 'txt':
        text = extract_text_from_txt(doc_path)
    elif file_type == 'html':
        text = extract_text_from_html(doc_path)
    elif file_type == 'pptx':
        text = extract_text_from_pptx(doc_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    # Generate output file path
    output_file_path = os.path.join('static/audio', f'speech_{selected_voice}.mp3')
    
    # Convert text to speech
    await text_to_speech_edge_tts(text, output_file_path, selected_voice)

    # ✅ Save generated audio file in `document_audio` table
    save_document_audio(document_id, selected_voice, output_file_path)

    return output_file_path
